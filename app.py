import streamlit as st
import requests
import datetime
import base64
import re
import urllib.parse
from html import escape
from io import BytesIO
import textwrap



# =========================
# Page config
# =========================
st.set_page_config(
    page_title="Signal Studio｜产品研判平台",
    page_icon="◌",
    layout="wide",
    initial_sidebar_state="collapsed",
)


# =========================
# Light, clean UI
# =========================
def inject_css():
    st.markdown(
        """
        <style>
        :root {
            --bg: #f8fbff;
            --ink: #162033;
            --muted: #66748a;
            --blue: #5d78ff;
            --blue-deep: #23366f;
            --blue-soft: #edf2ff;
            --warm: #b88758;
            --line: rgba(93,120,255,0.14);
            --panel: rgba(255,255,255,0.88);
            --shadow: 0 22px 58px rgba(22, 32, 51, 0.08);
        }

        .stApp {
            background:
                radial-gradient(circle at 12% 8%, rgba(93,120,255,0.10), transparent 24%),
                radial-gradient(circle at 94% 6%, rgba(184,135,88,0.08), transparent 24%),
                linear-gradient(180deg, #fbfdff 0%, #f3f7fd 100%);
            color: var(--ink);
        }

        .block-container {
            max-width: 1260px;
            padding-top: 1.4rem;
            padding-bottom: 4rem;
        }

        header[data-testid="stHeader"] {
            background: rgba(251,253,255,0.88);
            backdrop-filter: blur(14px);
            border-bottom: 1px solid rgba(93,120,255,0.06);
        }

        section[data-testid="stSidebar"] {
            background: linear-gradient(180deg, rgba(248,251,255,0.98) 0%, rgba(238,244,255,0.98) 100%);
            border-right: 1px solid rgba(93,120,255,0.12);
        }
        section[data-testid="stSidebar"] * { color: var(--ink) !important; }
        section[data-testid="stSidebar"] label,
        section[data-testid="stSidebar"] p,
        section[data-testid="stSidebar"] .stCaption {
            color: var(--muted) !important;
        }
        section[data-testid="stSidebar"] input,
        section[data-testid="stSidebar"] textarea,
        section[data-testid="stSidebar"] [data-baseweb="select"] > div,
        section[data-testid="stSidebar"] [data-baseweb="base-input"] {
            background: rgba(255,255,255,0.96) !important;
            border: 1px solid rgba(93,120,255,0.16) !important;
            border-radius: 16px !important;
            color: var(--ink) !important;
            box-shadow: none !important;
        }
        input::placeholder, textarea::placeholder {
            color: #95a1b5 !important;
            opacity: 1 !important;
        }

        h1, h2, h3, h4 { letter-spacing: -0.03em; }

        .top-hero {
            padding: 28px 30px;
            border-radius: 28px;
            background: linear-gradient(135deg, rgba(255,255,255,0.94) 0%, rgba(240,246,255,0.96) 100%);
            border: 1px solid rgba(93,120,255,0.12);
            box-shadow: var(--shadow);
            margin-bottom: 18px;
        }
        .eyebrow {
            font-size: 12px;
            letter-spacing: 0.28em;
            color: var(--blue);
            font-weight: 800;
            text-transform: uppercase;
            margin-bottom: 8px;
        }
        .hero-title {
            font-size: clamp(34px, 4.8vw, 58px);
            line-height: 1;
            color: var(--blue-deep);
            font-weight: 760;
            margin-bottom: 10px;
        }
        .hero-title span { color: var(--blue); font-weight: 520; }
        .hero-desc {
            max-width: 760px;
            color: var(--muted);
            font-size: 14px;
            line-height: 1.78;
            margin: 0;
        }
        .mini-steps {
            display: flex;
            flex-wrap: wrap;
            gap: 8px;
            margin-top: 16px;
        }
        .mini-chip {
            padding: 8px 12px;
            background: rgba(255,255,255,0.72);
            border: 1px solid rgba(93,120,255,0.10);
            border-radius: 999px;
            color: var(--blue-deep);
            font-size: 12px;
            font-weight: 650;
        }

        .query-preview {
            padding: 10px 12px;
            background: var(--blue-soft);
            border: 1px solid rgba(93,120,255,0.12);
            border-radius: 16px;
            color: var(--blue-deep);
            font-size: 13px;
            margin: 8px 0 12px;
        }

        .stTextInput label, .stTextArea label, .stSelectbox label, .stMultiSelect label, .stRadio label {
            color: var(--ink) !important;
            font-weight: 650;
        }
        .stTextInput input,
        .stTextArea textarea,
        [data-baseweb="select"] > div,
        [data-baseweb="base-input"] {
            background: rgba(255,255,255,0.90) !important;
            border: 1px solid rgba(93,120,255,0.12) !important;
            border-radius: 16px !important;
            color: var(--ink) !important;
            box-shadow: none !important;
        }
        .stTextInput input:focus, .stTextArea textarea:focus {
            border-color: rgba(93,120,255,0.32) !important;
            box-shadow: 0 0 0 4px rgba(93,120,255,0.08) !important;
        }

        .stButton > button, .stDownloadButton > button {
            min-height: 46px;
            border-radius: 999px !important;
            border: 1px solid rgba(93,120,255,0.16) !important;
            background: linear-gradient(135deg, #4968ee 0%, #7087ff 100%) !important;
            color: #ffffff !important;
            font-size: 13px;
            font-weight: 760;
            letter-spacing: 0.04em;
            box-shadow: 0 14px 30px rgba(73,104,238,0.16);
            transition: all .18s ease;
        }
        .stButton > button:hover, .stDownloadButton > button:hover {
            transform: translateY(-1px);
            box-shadow: 0 20px 36px rgba(73,104,238,0.22);
        }

        div[data-testid="stLinkButton"] a {
            border-radius: 18px !important;
            border: 1px solid rgba(93,120,255,0.12) !important;
            background: rgba(255,255,255,0.92) !important;
            color: var(--blue-deep) !important;
            min-height: 48px !important;
            box-shadow: 0 12px 24px rgba(22,32,51,0.05) !important;
            font-weight: 700 !important;
        }
        div[data-testid="stLinkButton"] a:hover {
            border-color: rgba(93,120,255,0.28) !important;
            transform: translateY(-1px);
        }

        div[data-testid="stAlert"] {
            background: rgba(255,255,255,0.78);
            border: 1px solid rgba(93,120,255,0.10);
            border-radius: 18px;
        }
        div[data-testid="stExpander"] {
            background: rgba(255,255,255,0.74) !important;
            border: 1px solid rgba(93,120,255,0.10) !important;
            border-radius: 18px !important;
        }

        .download-link {
            display: inline-block;
            text-decoration: none !important;
            padding: 12px 18px;
            margin-top: 12px;
            border-radius: 999px;
            background: linear-gradient(135deg, #4968ee 0%, #7087ff 100%);
            color: white !important;
            font-size: 13px;
            font-weight: 760;
        }
        .hint-box {
            background: rgba(255,255,255,0.78);
            border: 1px solid rgba(93,120,255,0.10);
            border-radius: 18px;
            padding: 14px 16px;
            color: var(--muted);
            line-height: 1.75;
            font-size: 13px;
        }
        .micro-tip { color: var(--muted); font-size: 12px; line-height: 1.7; }

        .report-shell {
            background: rgba(255,255,255,0.94);
            border: 1px solid rgba(93,120,255,0.10);
            border-radius: 26px;
            box-shadow: 0 28px 64px rgba(22,32,51,0.08);
            overflow: hidden;
            margin-top: 6px;
        }
        .report-head {
            padding: 26px 28px 22px;
            border-bottom: 1px solid rgba(93,120,255,0.10);
            background: linear-gradient(180deg, rgba(245,248,255,0.95) 0%, rgba(255,255,255,0.92) 100%);
        }
        .report-kicker {
            font-size: 11px;
            text-transform: uppercase;
            letter-spacing: 0.22em;
            color: #5d78ff;
            font-weight: 800;
            margin-bottom: 10px;
        }
        .report-title {
            font-size: 30px;
            line-height: 1.15;
            font-weight: 760;
            color: #162033;
            margin: 0 0 8px;
        }
        .report-meta {
            color: #7a879b;
            font-size: 12px;
        }
        .summary-grid {
            display: grid;
            grid-template-columns: repeat(auto-fit, minmax(210px, 1fr));
            gap: 12px;
            padding: 18px 28px 6px;
        }
        .summary-card {
            padding: 16px 16px 14px;
            background: linear-gradient(180deg, #f7faff 0%, #ffffff 100%);
            border: 1px solid rgba(93,120,255,0.10);
            border-radius: 18px;
        }
        .summary-card .num {
            font-size: 11px;
            color: #5d78ff;
            font-weight: 800;
            letter-spacing: 0.12em;
            text-transform: uppercase;
            margin-bottom: 8px;
        }
        .summary-card p {
            margin: 0;
            font-size: 14px;
            line-height: 1.7;
            color: #162033;
        }
        .report-body {
            padding: 8px 28px 28px;
        }
        .report-section {
            margin-top: 18px;
            padding-top: 18px;
            border-top: 1px solid rgba(93,120,255,0.08);
        }
        .report-section h2 {
            margin: 0 0 12px;
            font-size: 19px;
            line-height: 1.35;
            color: #23366f;
            font-weight: 760;
        }
        .report-section p {
            margin: 8px 0;
            line-height: 1.78;
            font-size: 14px;
            color: #22314f;
        }
        .report-list {
            margin: 0;
            padding-left: 1.2rem;
        }
        .report-list li {
            margin: 8px 0;
            line-height: 1.76;
            color: #22314f;
            font-size: 14px;
        }
        .report-quote {
            background: #fbf8f3;
            border: 1px solid rgba(184,135,88,0.16);
            border-left: 3px solid #b88758;
            border-radius: 14px;
            padding: 12px 14px;
            margin: 10px 0;
            color: #32405c;
            font-size: 14px;
            line-height: 1.76;
        }
        .table-wrap {
            overflow-x: auto;
            border: 1px solid rgba(93,120,255,0.10);
            border-radius: 18px;
            margin-top: 12px;
        }
        .report-table {
            width: 100%;
            border-collapse: collapse;
            font-size: 13px;
            background: #fff;
        }
        .report-table th,
        .report-table td {
            padding: 12px 12px;
            text-align: left;
            vertical-align: top;
            border-bottom: 1px solid rgba(93,120,255,0.08);
            border-right: 1px solid rgba(93,120,255,0.06);
            line-height: 1.72;
            color: #22314f;
        }
        .report-table th {
            background: #eff4ff;
            color: #162033;
            font-weight: 760;
        }
        .report-table tr:last-child td { border-bottom: none; }
        .report-table th:last-child, .report-table td:last-child { border-right: none; }
        .report-footer {
            padding: 18px 28px 24px;
            font-size: 11px;
            color: #7a879b;
            border-top: 1px solid rgba(93,120,255,0.08);
            letter-spacing: 0.08em;
        }

        @media (max-width: 800px) {
            .top-hero { padding: 22px; }
            .hero-title { font-size: 38px; }
            .report-title { font-size: 24px; }
            .summary-grid, .report-body, .report-head, .report-footer { padding-left: 18px; padding-right: 18px; }
        }
        </style>
        """,
        unsafe_allow_html=True,
    )


inject_css()


# =========================
# Search definitions
# =========================
INTENT_SUFFIX = {
    "全网口碑": "口碑 体验 测评",
    "吐槽痛点": "吐槽 缺点 问题 避雷",
    "测评体验": "测评 体验 上手 深度评测",
    "争议翻车": "争议 翻车 差评 售后",
    "购买建议": "值得买吗 购买建议 真实体验",
    "价格价值感": "价格 值不值得 性价比 贵不贵",
    "竞品对比": "对比 vs 区别 怎么选",
}

SEARCH_ENGINES = [
    {"group": "UGC 社媒", "name": "小红书", "url": "https://www.xiaohongshu.com/search_result?keyword={q}"},
    {"group": "UGC 社媒", "name": "微博", "url": "https://s.weibo.com/weibo?q={q}"},
    {"group": "UGC 社媒", "name": "知乎", "url": "https://www.zhihu.com/search?type=content&q={q}"},
    {"group": "UGC 社媒", "name": "微信文章", "url": "https://weixin.sogou.com/weixin?type=2&query={q}"},
    {"group": "视频测评", "name": "抖音", "url": "https://www.douyin.com/search/{q}"},
    {"group": "视频测评", "name": "B站", "url": "https://search.bilibili.com/all?keyword={q}"},
    {"group": "视频测评", "name": "快手", "url": "https://www.kuaishou.com/search/video?searchKey={q}"},
    {"group": "搜索引擎", "name": "百度", "url": "https://www.baidu.com/s?wd={q}"},
    {"group": "搜索引擎", "name": "Google", "url": "https://www.google.com/search?q={q}"},
    {"group": "搜索引擎", "name": "必应", "url": "https://www.bing.com/search?q={q}"},
    {"group": "电商口碑", "name": "京东", "url": "https://search.jd.com/Search?keyword={q}"},
    {"group": "电商口碑", "name": "什么值得买", "url": "https://search.smzdm.com/?c=home&s={q}"},
]


API_PRESETS = {
    "阿里云 DashScope": {
        "base": "https://dashscope.aliyuncs.com/compatible-mode/v1",
        "models": ["deepseek-v3", "qwen-max", "qwen-plus", "qwen-turbo", "qwen-vl-plus", "qwen-vl-max"],
    },
    "OpenRouter": {
        "base": "https://openrouter.ai/api/v1",
        "models": ["deepseek/deepseek-chat-v3-0324:free", "openai/gpt-4.1-mini", "openai/gpt-4o-mini", "google/gemini-2.5-flash-preview", "google/gemini-2.0-flash-001"],
    },
    "SiliconFlow": {
        "base": "https://api.siliconflow.cn/v1",
        "models": ["deepseek-ai/DeepSeek-V3", "Qwen/Qwen2.5-72B-Instruct", "Qwen/Qwen2.5-VL-72B-Instruct", "THUDM/GLM-4-9B-Chat"],
    },
}


def compose_query(product: str, intent: str, aliases: str = "", competitor: str = "") -> str:
    parts = [product.strip()]
    if competitor.strip():
        parts.append(competitor.strip())
    if aliases.strip():
        parts.append(aliases.strip())
    parts.append(INTENT_SUFFIX.get(intent, "口碑 体验 测评"))
    return " ".join([p for p in parts if p]).strip()


def build_search_links(query: str, groups: list[str] | None = None):
    links = []
    for engine in SEARCH_ENGINES:
        if groups and engine["group"] not in groups:
            continue
        encoded = urllib.parse.quote(query)
        links.append({
            "group": engine["group"],
            "name": engine["name"],
            "url": engine["url"].format(q=encoded),
        })
    return links


def sanitize_filename(name: str) -> str:
    name = re.sub(r"[\\/:*?\"<>|\s]+", "_", name.strip())
    return name[:80] or "product_report"


def is_ready_to_analyze(api_key: str, *texts: str) -> bool:
    if not api_key.strip():
        st.warning("请先在侧边配置栏输入 API Key。")
        return False
    if any(not text.strip() for text in texts):
        st.warning("请先填写产品名，并粘贴真实语料。")
        return False
    return True


def is_ready_to_analyze_multimodal(api_key: str, required_names: list[str], evidence_groups: list[tuple[str, str, list]]) -> bool:
    if not api_key.strip():
        st.warning("请先在侧边配置栏输入 API Key。")
        return False
    if any(not name.strip() for name in required_names):
        st.warning("请先填写产品名/竞品名。")
        return False
    missing = []
    for label, text_value, image_files in evidence_groups:
        has_text = bool((text_value or "").strip())
        has_image = bool(image_files)
        if not has_text and not has_image:
            missing.append(label)
    if missing:
        st.warning("请为以下对象粘贴语料或上传图片证据：" + "、".join(missing))
        return False
    return True


def get_corpus_template(product_name: str = "产品名") -> str:
    product_name = product_name.strip() or "产品名"
    return f"""【产品】{product_name}
【采集日期】{datetime.date.today()}
【平台】小红书 / 微博 / 抖音 / B站 / 京东 / 其他
【搜索词】{product_name} 体验 吐槽 测评

【原始语料】
1. 平台：
   用户原话/评论摘录：
   互动情况：点赞/评论/转发，如无可不填
   初步标签：价格 / 影像 / 外观 / 系统 / 售后 / 购买建议

2. 平台：
   用户原话/评论摘录：
   互动情况：
   初步标签：

【备注】
- 尽量保留用户原话，不要提前改写。
- 不确定的信息不要补充，让模型在报告中写“暂无提及”。"""


# =========================
# Rendering helpers
# =========================
def render_header():
    st.markdown(
        """
        <section class="top-hero">
            <div class="eyebrow">SIGNAL STUDIO · PRODUCT INTELLIGENCE</div>
            <div class="hero-title">Strategy <span>Signal Deck</span></div>
            <p class="hero-desc">
            轻量产品情报台：先用便携搜索坞快速取证，再粘贴真实语料生成单品研判或竞品攻防。
            这一版重点优化了报告产出：页面更克制、版式更干净、内容更适合直接给老板看。
            </p>
            <div class="mini-steps">
                <span class="mini-chip">01 搜索取证</span>
                <span class="mini-chip">02 粘贴语料</span>
                <span class="mini-chip">03 生成老板可读简报</span>
            </div>
        </section>
        """,
        unsafe_allow_html=True,
    )


def render_link_buttons(links: list[dict], columns_per_row: int = 4):
    for start in range(0, len(links), columns_per_row):
        row = links[start:start + columns_per_row]
        cols = st.columns(columns_per_row)
        for idx, col in enumerate(cols):
            if idx < len(row):
                item = row[idx]
                col.link_button(
                    f"{item['name']} · {item['group']}",
                    item["url"],
                    use_container_width=True,
                )


def render_search_dock():
    if "quick_product" not in st.session_state:
        st.session_state.quick_product = ""
    if "single_product_name" not in st.session_state:
        st.session_state.single_product_name = ""
    if "main_product" not in st.session_state:
        st.session_state.main_product = ""
    if "competitor_product" not in st.session_state:
        st.session_state.competitor_product = ""

    with st.container(border=True):
        st.markdown("### 便携搜索坞")
        st.caption("输入产品名后，直接打开各平台搜索入口。这里全部改成原生按钮，不再用 HTML 卡片，避免页面把代码渲染出来。")

        c1, c2, c3, c4 = st.columns([1.35, 0.9, 1.0, 0.95])
        with c1:
            product = st.text_input("产品关键词", key="quick_product", placeholder="例如：vivo X Fold5 / OPPO Find N6")
        with c2:
            intent = st.selectbox("搜索意图", list(INTENT_SUFFIX.keys()), index=0)
        with c3:
            aliases = st.text_input("补充词 / 别名", placeholder="例如：折叠屏 蓝厂")
        with c4:
            competitor = st.text_input("竞品词，可空", placeholder="例如：Find N6")

        query = compose_query(product, intent, aliases, competitor)

        if product.strip():
            st.markdown(f"<div class='query-preview'>当前搜索式：{escape(query)}</div>", unsafe_allow_html=True)

            sync_cols = st.columns(3)
            if sync_cols[0].button("同步到单品分析", use_container_width=True):
                st.session_state.single_product_name = product.strip()
                st.rerun()
            if sync_cols[1].button("同步为本品", use_container_width=True):
                st.session_state.main_product = product.strip()
                st.rerun()
            if sync_cols[2].button("同步为竞品", use_container_width=True):
                st.session_state.competitor_product = product.strip()
                st.rerun()

            group_tabs = st.tabs(["常用入口", "UGC 社媒", "视频测评", "搜索引擎", "电商口碑"])
            with group_tabs[0]:
                common = build_search_links(query, ["UGC 社媒", "视频测评", "搜索引擎"])[:8]
                render_link_buttons(common, columns_per_row=4)
            with group_tabs[1]:
                render_link_buttons(build_search_links(query, ["UGC 社媒"]), columns_per_row=4)
            with group_tabs[2]:
                render_link_buttons(build_search_links(query, ["视频测评"]), columns_per_row=3)
            with group_tabs[3]:
                render_link_buttons(build_search_links(query, ["搜索引擎"]), columns_per_row=3)
            with group_tabs[4]:
                render_link_buttons(build_search_links(query, ["电商口碑"]), columns_per_row=2)

            with st.expander("复制搜索式 / URL 清单"):
                st.code(query, language="text")
                url_lines = [f"{item['name']}：{item['url']}" for item in build_search_links(query)]
                st.code("\n".join(url_lines), language="text")
        else:
            st.info("先输入产品关键词，系统会生成小红书、微博、抖音、B站、搜索引擎、电商评价入口。")


def render_evidence_helper(product_name: str):
    with st.expander("语料整理模板 / 证据质量提示", expanded=False):
        st.markdown(
            """
            <div class="hint-box">
            建议至少收集 3 类语料：社媒真实吐槽、媒体/博主测评摘录、电商评价或购买问答。
            分析前不要把评论过度润色，保留原话更容易生成有效洞察。
            如果是给老板看，建议优先贴“最能代表主流情绪”的评论，而不是堆很多重复句子。
            </div>
            """,
            unsafe_allow_html=True,
        )
        st.code(get_corpus_template(product_name), language="text")


def render_uploaded_images(files, caption_prefix="图片证据"):
    if not files:
        return
    st.caption(f"已上传 {len(files)} 张{caption_prefix}。图片会随报告一起送入支持视觉能力的模型分析。")
    cols = st.columns(min(4, max(1, len(files))))
    for idx, file in enumerate(files[:8]):
        with cols[idx % len(cols)]:
            st.image(file, caption=file.name, use_container_width=True)



# =========================
# LLM engine and prompt templates
# =========================
def image_file_to_block(uploaded_file):
    data = uploaded_file.getvalue()
    mime = uploaded_file.type or "image/png"
    b64 = base64.b64encode(data).decode("utf-8")
    return {
        "type": "image_url",
        "image_url": {"url": f"data:{mime};base64,{b64}"},
    }


def build_user_content(prompt: str, image_groups: list[tuple[str, list]] | None = None):
    if not image_groups:
        return prompt
    content = [{"type": "text", "text": prompt}]
    for label, files in image_groups:
        if not files:
            continue
        content.append({"type": "text", "text": f"以下图片证据归属：{label}。请只读取图片中可见的信息，不要推断图片之外的事实。"})
        for file in files:
            content.append(image_file_to_block(file))
    return content


def analyze_with_llm(prompt: str, api_key: str, model_name: str, api_base: str, image_groups: list[tuple[str, list]] | None = None):
    url = api_base.rstrip("/") + "/chat/completions"
    headers = {"Authorization": "Bearer " + api_key.strip(), "Content-Type": "application/json"}
    payload = {
        "model": model_name.strip(),
        "messages": [
            {
                "role": "system",
                "content": (
                    "你是一名严谨的产品策略与用户舆情分析师，服务对象是业务负责人、产品负责人和市场负责人。"
                    "你只能根据用户提供的真实语料和图片证据进行分析，严禁编造参数、销量、配置、发布时间、品牌动作或用户评价。"
                    "如果语料或图片中没有提到某个信息，必须写‘暂无提及’。"
                    "如用户上传截图、海报、评论区截图、表格或产品图片，你可以读取其中可见文字、画面元素和明确可见的信息，但不得推断图片之外的信息。"
                    "输出必须是老板可直接阅读的内部简报风格：少废话、少空话、信息密度高、判断清晰、动作明确。"
                    "不要写长篇铺垫、方法说明、免责声明、过程性废话。"
                    "每个核心判断尽量给出对应证据。表格请使用 Markdown 表格。"
                ),
            },
            {"role": "user", "content": build_user_content(prompt, image_groups)},
        ],
        "temperature": 0.0,
    }
    try:
        res = requests.post(url, headers=headers, json=payload, timeout=90)
        res.raise_for_status()
        data = res.json()
        return data["choices"][0]["message"]["content"]
    except Exception as e:
        return f"ERROR: 引擎响应异常，请检查 API Key、模型名称、接口地址，或确认当前模型是否支持图片输入。错误信息：{e}"


def build_single_prompt(product_name: str, product_type: str, focus: list[str], corpus: str) -> str:
    focus_text = "、".join(focus) if focus else "用户体验、核心卖点、负面反馈、购买阻力、传播机会"
    return f"""
请仅根据下方语料，撰写《{product_name} 单品研判简报》。

产品名称：{product_name}
产品类型：{product_type}
重点关注：{focus_text}

请严格遵守以下要求：
- 只基于真实语料判断，严禁补充外部常识、产品参数、销量、发布时间、品牌策略等未在语料出现的信息。
- 这份报告默认直接给老板看，必须克制、简洁、专业，不要写空话套话。
- 不要复述任务，不要写“以下是报告”“基于以上分析”等废话。
- 每一条结论尽量给出证据摘录。
- 如果上传了图片证据，请读取图片中可见的文字、截图内容、画面元素，并在证据摘录中标注“图片证据”。
- 如果某部分证据不足，统一简洁写“暂无提及”，不要重复解释。

输出结构请固定如下：

一、老板先看
- 用 3—5 条要点概括当前最重要的判断。
- 每条尽量包含：现象 / 判断 / 对业务的意义。

二、用户反馈总览
请输出 Markdown 表格，字段为：
| 主题 | 情绪倾向 | 用户反馈要点 | 证据摘录 | 对业务的影响 |
主题可自适应，如外观设计、影像/性能、系统体验、价格价值感、服务/渠道、品牌认知、购买决策等。

三、关键风险
请输出 Markdown 表格，字段为：
| 风险点 | 具体表现 | 证据摘录 | 风险等级 | 建议动作 |
只列真正值得管理层关注的点，不要凑数。

四、可放大的卖点 / 机会
- 提炼 3—5 条有证据支撑的机会点。
- 说明为什么它值得放大。

五、建议动作
请分成三个小部分输出：
- 产品侧：短期该优化什么
- 营销侧：传播上该强调/规避什么
- 取证侧：下一轮还要补什么语料
每条建议都尽量与语料对应。

六、代表性用户原声
- 提取 3—5 条最有代表性的用户原话或高还原度近似原话。
- 不要过多，重质不重量。

七、信息缺口
- 用 2—4 条列出当前还无法判断、但决策上可能重要的信息。

真实语料如下：
{corpus}
""".strip()


def build_compare_prompt(main_product: str, competitor_product: str, product_type: str, focus: list[str], main_corpus: str, competitor_corpus: str) -> str:
    focus_text = "、".join(focus) if focus else "用户体验、核心卖点、负面反馈、购买阻力、传播机会"
    return f"""
请仅根据下方两组语料，撰写《{main_product} vs {competitor_product} 竞品攻防简报》。

本品：{main_product}
竞品：{competitor_product}
产品类型：{product_type}
重点关注：{focus_text}

请严格遵守以下要求：
- 只根据用户提供语料判断，严禁补充外部参数、销量、发布时间、品牌动作或行业常识。
- 这份内容默认直接给老板看，要像内部策略简报：结论明确、语言克制、可直接用于讨论。
- 不要写长篇铺垫，不要凑字数，不要强行制造差异。
- 若某一方证据不足，简洁写“暂无提及”。
- 每个对比结论尽量同时体现本品与竞品证据。
- 如果上传了图片证据，请区分本品图片与竞品图片，只读取可见信息，并在证据摘录中标注“图片证据”。

输出结构请固定如下：

一、老板先看
- 用 3—5 条说清楚本品当前最值得管理层关注的攻防结论。

二、核心对比总览
请输出 Markdown 表格，字段为：
| 对比维度 | {main_product} 用户反馈 | {competitor_product} 用户反馈 | 结论判断 | 对业务的启示 |
维度根据语料自适应，如外观设计、性能、系统、价格、品牌、高端认知、渠道服务等。

三、本品风险点（需要防守）
请输出 Markdown 表格，字段为：
| 风险点 | 具体表现 | 证据摘录 | 风险等级 | 防守动作 |

四、本品机会点（可以进攻）
- 提炼 3—5 条本品可主动进攻的点。
- 每条说明对应依据与可用场景。

五、建议动作
请分成三个小部分输出：
- 产品侧：需优先补的短板
- 营销侧：传播上应强调/回避的点
- 取证侧：还应补哪些平台和关键词

六、双方代表性用户原声
- 本品列 2—3 条，竞品列 2—3 条。
- 只保留最能说明问题的内容。

七、信息缺口
- 用 2—4 条列出当前证据不足、但会影响判断的内容。

【本品真实语料】
{main_corpus}

【竞品真实语料】
{competitor_corpus}
""".strip()


# =========================
# Report formatting helpers
# =========================
def normalize_line(line: str) -> str:
    line = line.strip()
    line = re.sub(r"^#+\s*", "", line)
    line = line.replace("**", "")
    line = line.replace("__", "")
    return line.strip()


def is_section_heading(line: str) -> bool:
    line = normalize_line(line)
    return bool(re.match(r"^[一二三四五六七八九十]+[、.．]\s*", line))


def parse_sections(text_content: str):
    sections = []
    current_title = ""
    current_lines: list[str] = []
    for raw in text_content.splitlines():
        line = normalize_line(raw)
        if not line:
            current_lines.append("")
            continue
        if is_section_heading(line):
            if current_title or current_lines:
                sections.append((current_title, current_lines))
            current_title = line
            current_lines = []
        else:
            current_lines.append(line)
    if current_title or current_lines:
        sections.append((current_title, current_lines))
    return sections


def detect_blocks(lines: list[str]):
    cleaned = [normalize_line(l) for l in lines]
    blocks = []
    i = 0
    while i < len(cleaned):
        line = cleaned[i].strip()
        if not line:
            i += 1
            continue

        if (
            "|" in line and i + 1 < len(cleaned)
            and "|" in cleaned[i + 1]
            and re.fullmatch(r"[|:\-\s]+", cleaned[i + 1]) is not None
        ):
            table_lines = [line, cleaned[i + 1].strip()]
            i += 2
            while i < len(cleaned) and "|" in cleaned[i].strip():
                table_lines.append(cleaned[i].strip())
                i += 1
            blocks.append(("table", table_lines))
            continue

        if re.match(r"^[-•—*]\s+", line):
            items = []
            while i < len(cleaned) and re.match(r"^[-•—*]\s+", cleaned[i].strip()):
                items.append(re.sub(r"^[-•—*]\s+", "", cleaned[i].strip()))
                i += 1
            blocks.append(("ul", items))
            continue

        if re.match(r"^\d+[、.．]\s*", line):
            items = []
            while i < len(cleaned) and re.match(r"^\d+[、.．]\s*", cleaned[i].strip()):
                items.append(re.sub(r"^\d+[、.．]\s*", "", cleaned[i].strip()))
                i += 1
            blocks.append(("ol", items))
            continue

        paras = []
        while i < len(cleaned):
            probe = cleaned[i].strip()
            if not probe:
                if paras:
                    i += 1
                    break
                i += 1
                continue
            if (
                re.match(r"^[-•—*]\s+", probe)
                or re.match(r"^\d+[、.．]\s*", probe)
                or (
                    "|" in probe and i + 1 < len(cleaned)
                    and "|" in cleaned[i + 1]
                    and re.fullmatch(r"[|:\-\s]+", cleaned[i + 1]) is not None
                )
            ):
                break
            paras.append(probe)
            i += 1
        if paras:
            blocks.append(("p", paras))
    return blocks


def parse_markdown_table(table_lines: list[str]):
    rows = []
    for idx, line in enumerate(table_lines):
        if idx == 1 and re.fullmatch(r"[|:\-\s]+", line):
            continue
        parts = [normalize_line(c) for c in line.split("|")]
        parts = [p for p in parts if p != ""]
        if parts:
            rows.append(parts)
    if not rows:
        return [], []
    header = rows[0]
    body = rows[1:]
    return header, body


def render_table_html(table_lines: list[str]) -> str:
    header, body = parse_markdown_table(table_lines)
    if not header:
        return ""
    ths = "".join(f"<th>{escape(h)}</th>" for h in header)
    body_rows = ""
    for row in body:
        if len(row) < len(header):
            row = row + [""] * (len(header) - len(row))
        tds = "".join(f"<td>{escape(cell)}</td>" for cell in row[:len(header)])
        body_rows += f"<tr>{tds}</tr>"
    return f"<div class='table-wrap'><table class='report-table'><thead><tr>{ths}</tr></thead><tbody>{body_rows}</tbody></table></div>"


def render_blocks_html(title: str, lines: list[str]) -> str:
    title_clean = normalize_line(title)
    quote_mode = any(key in title_clean for key in ["原声", "用户原声", "代表性"])
    html_parts = []
    for block_type, content in detect_blocks(lines):
        if block_type == "table":
            html_parts.append(render_table_html(content))
        elif block_type == "ul":
            if quote_mode:
                for item in content:
                    html_parts.append(f"<div class='report-quote'>“{escape(item)}”</div>")
            else:
                items = "".join(f"<li>{escape(item)}</li>" for item in content)
                html_parts.append(f"<ul class='report-list'>{items}</ul>")
        elif block_type == "ol":
            items = "".join(f"<li>{escape(item)}</li>" for item in content)
            html_parts.append(f"<ol class='report-list'>{items}</ol>")
        elif block_type == "p":
            for para in content:
                if quote_mode and para and not any(tag in para for tag in ["：", ":"]) and len(para) < 80:
                    html_parts.append(f"<div class='report-quote'>“{escape(para)}”</div>")
                else:
                    html_parts.append(f"<p>{escape(para)}</p>")
    return "".join(html_parts)


def extract_summary_points(sections: list[tuple[str, list[str]]], max_items: int = 4):
    if not sections:
        return []
    first_title, first_lines = sections[0]
    title_text = normalize_line(first_title)
    if not any(k in title_text for k in ["老板先看", "核心结论", "对比结论"]):
        return []
    points = []
    for block_type, content in detect_blocks(first_lines):
        if block_type in {"ul", "ol"}:
            points.extend(content)
        elif block_type == "p":
            points.extend(content)
        if len(points) >= max_items:
            break
    return [p for p in points[:max_items] if p]


def build_report_body_html(text_content: str, title: str) -> str:
    sections = parse_sections(text_content)
    summary_points = extract_summary_points(sections)

    summary_html = ""
    if summary_points:
        cards = []
        for idx, point in enumerate(summary_points, start=1):
            cards.append(
                f"<div class='summary-card'><div class='num'>Key {idx:02d}</div><p>{escape(point)}</p></div>"
            )
        summary_html = f"<div class='summary-grid'>{''.join(cards)}</div>"

    body_sections = []
    for idx, (section_title, section_lines) in enumerate(sections):
        title_clean = normalize_line(section_title) or f"部分 {idx+1}"
        if idx == 0 and summary_points:
            section_html = render_blocks_html(section_title, section_lines)
            # keep section but summary already highlighted visually
            body_sections.append(
                f"<section class='report-section'><h2>{escape(title_clean)}</h2>{section_html}</section>"
            )
        else:
            section_html = render_blocks_html(section_title, section_lines)
            body_sections.append(
                f"<section class='report-section'><h2>{escape(title_clean)}</h2>{section_html}</section>"
            )

    if not body_sections:
        body_sections = [f"<section class='report-section'><p>{escape(text_content)}</p></section>"]

    meta = f"生成日期：{datetime.date.today()} · 内部参考简报"
    return (
        f"<div class='report-shell'>"
        f"<div class='report-head'><div class='report-kicker'>SIGNAL STUDIO · EXECUTIVE BRIEF</div>"
        f"<div class='report-title'>{escape(title)}</div><div class='report-meta'>{meta}</div></div>"
        f"{summary_html}"
        f"<div class='report-body'>{''.join(body_sections)}</div>"
        f"<div class='report-footer'>EVIDENCE-BASED · CONCISE · ACTIONABLE</div>"
        f"</div>"
    )


def build_download_html(text_content: str, title: str) -> str:
    body = build_report_body_html(text_content, title)
    return f"""
    <html>
    <head>
        <meta charset="utf-8" />
        <meta name="viewport" content="width=device-width, initial-scale=1.0" />
        <title>{escape(title)}</title>
        <style>
            body {{
                margin: 0;
                padding: 28px;
                background: linear-gradient(180deg, #f8fbff 0%, #eef4ff 100%);
                font-family: -apple-system, BlinkMacSystemFont, 'PingFang SC', 'Microsoft YaHei', sans-serif;
            }}
            .report-shell {{ background: rgba(255,255,255,0.96); border: 1px solid rgba(93,120,255,0.10); border-radius: 24px; box-shadow: 0 30px 80px rgba(22,32,51,0.10); overflow: hidden; max-width: 1080px; margin: 0 auto; }}
            .report-head {{ padding: 28px 30px 22px; border-bottom: 1px solid rgba(93,120,255,0.10); background: linear-gradient(180deg, rgba(245,248,255,0.95) 0%, rgba(255,255,255,0.92) 100%); }}
            .report-kicker {{ font-size: 11px; text-transform: uppercase; letter-spacing: 0.22em; color: #5d78ff; font-weight: 800; margin-bottom: 10px; }}
            .report-title {{ font-size: 32px; line-height: 1.15; font-weight: 760; color: #162033; margin: 0 0 8px; }}
            .report-meta {{ color: #7a879b; font-size: 12px; }}
            .summary-grid {{ display: grid; grid-template-columns: repeat(auto-fit, minmax(210px, 1fr)); gap: 12px; padding: 18px 30px 8px; }}
            .summary-card {{ padding: 16px 16px 14px; background: linear-gradient(180deg, #f7faff 0%, #ffffff 100%); border: 1px solid rgba(93,120,255,0.10); border-radius: 18px; }}
            .summary-card .num {{ font-size: 11px; color: #5d78ff; font-weight: 800; letter-spacing: 0.12em; text-transform: uppercase; margin-bottom: 8px; }}
            .summary-card p {{ margin: 0; font-size: 14px; line-height: 1.7; color: #162033; }}
            .report-body {{ padding: 8px 30px 30px; }}
            .report-section {{ margin-top: 18px; padding-top: 18px; border-top: 1px solid rgba(93,120,255,0.08); }}
            .report-section h2 {{ margin: 0 0 12px; font-size: 20px; line-height: 1.35; color: #23366f; font-weight: 760; }}
            .report-section p {{ margin: 8px 0; line-height: 1.8; font-size: 14px; color: #22314f; }}
            .report-list {{ margin: 0; padding-left: 1.25rem; }}
            .report-list li {{ margin: 8px 0; line-height: 1.76; color: #22314f; font-size: 14px; }}
            .report-quote {{ background: #fbf8f3; border: 1px solid rgba(184,135,88,0.16); border-left: 3px solid #b88758; border-radius: 14px; padding: 12px 14px; margin: 10px 0; color: #32405c; font-size: 14px; line-height: 1.76; }}
            .table-wrap {{ overflow-x: auto; border: 1px solid rgba(93,120,255,0.10); border-radius: 18px; margin-top: 12px; }}
            .report-table {{ width: 100%; border-collapse: collapse; font-size: 13px; background: #fff; }}
            .report-table th, .report-table td {{ padding: 12px 12px; text-align: left; vertical-align: top; border-bottom: 1px solid rgba(93,120,255,0.08); border-right: 1px solid rgba(93,120,255,0.06); line-height: 1.72; color: #22314f; }}
            .report-table th {{ background: #eff4ff; color: #162033; font-weight: 760; }}
            .report-table tr:last-child td {{ border-bottom: none; }}
            .report-table th:last-child, .report-table td:last-child {{ border-right: none; }}
            .report-footer {{ padding: 18px 30px 24px; font-size: 11px; color: #7a879b; border-top: 1px solid rgba(93,120,255,0.08); letter-spacing: 0.08em; }}
        </style>
    </head>
    <body>
        {body}
    </body>
    </html>
    """


def generate_html_report(text_content: str, title: str):
    html_doc = build_download_html(text_content, title)
    b64 = base64.b64encode(html_doc.encode("utf-8")).decode()
    filename = sanitize_filename(title)
    return f'<a class="download-link" href="data:text/html;base64,{b64}" download="{filename}.html">导出精美简报 · HTML/PDF</a>'


def clean_cell_text(value: str) -> str:
    return re.sub(r"<[^>]+>", "", normalize_line(value or ""))


def create_docx_report(report_text: str, title: str) -> bytes:
    from docx import Document
    from docx.shared import Pt, Inches
    from docx.enum.text import WD_ALIGN_PARAGRAPH

    doc = Document()
    sec = doc.sections[0]
    sec.top_margin = Inches(0.65)
    sec.bottom_margin = Inches(0.65)
    sec.left_margin = Inches(0.72)
    sec.right_margin = Inches(0.72)

    h = doc.add_heading(title, level=0)
    h.alignment = WD_ALIGN_PARAGRAPH.CENTER
    meta = doc.add_paragraph(f"生成日期：{datetime.date.today()} · 内部参考简报")
    meta.alignment = WD_ALIGN_PARAGRAPH.CENTER

    for section_title, lines in parse_sections(report_text):
        if section_title:
            doc.add_heading(normalize_line(section_title), level=1)
        for block_type, content in detect_blocks(lines):
            if block_type == "table":
                header, rows = parse_markdown_table(content)
                if not header:
                    continue
                table = doc.add_table(rows=1, cols=len(header))
                table.style = "Table Grid"
                for i, hcell in enumerate(header):
                    table.rows[0].cells[i].text = clean_cell_text(hcell)
                for row in rows:
                    cells = table.add_row().cells
                    row = row + [""] * (len(header) - len(row))
                    for i, cell in enumerate(row[:len(header)]):
                        cells[i].text = clean_cell_text(cell)
            elif block_type in {"ul", "ol"}:
                style = "List Bullet" if block_type == "ul" else "List Number"
                for item in content:
                    doc.add_paragraph(clean_cell_text(item), style=style)
            elif block_type == "p":
                for para in content:
                    doc.add_paragraph(clean_cell_text(para))

    for para in doc.paragraphs:
        for run in para.runs:
            run.font.name = "Microsoft YaHei"
            run.font.size = Pt(10.5)
    bio = BytesIO()
    doc.save(bio)
    return bio.getvalue()


def create_pdf_report(report_text: str, title: str) -> bytes:
    from reportlab.lib import colors
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.units import mm
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle
    from reportlab.pdfbase import pdfmetrics
    from reportlab.pdfbase.cidfonts import UnicodeCIDFont

    try:
        pdfmetrics.registerFont(UnicodeCIDFont("STSong-Light"))
        base_font = "STSong-Light"
    except Exception:
        base_font = "Helvetica"

    bio = BytesIO()
    doc = SimpleDocTemplate(bio, pagesize=A4, rightMargin=16*mm, leftMargin=16*mm, topMargin=16*mm, bottomMargin=16*mm)
    styles = getSampleStyleSheet()
    title_style = ParagraphStyle("CNTitle", parent=styles["Title"], fontName=base_font, fontSize=20, leading=26, textColor=colors.HexColor("#162033"), alignment=1, spaceAfter=10)
    meta_style = ParagraphStyle("CNMeta", parent=styles["Normal"], fontName=base_font, fontSize=8.5, leading=12, textColor=colors.HexColor("#7a879b"), alignment=1, spaceAfter=12)
    h_style = ParagraphStyle("CNHeading", parent=styles["Heading2"], fontName=base_font, fontSize=13.5, leading=18, textColor=colors.HexColor("#23366f"), spaceBefore=12, spaceAfter=8)
    p_style = ParagraphStyle("CNBody", parent=styles["BodyText"], fontName=base_font, fontSize=9.5, leading=15, textColor=colors.HexColor("#22314f"), spaceAfter=6)
    bullet_style = ParagraphStyle("CNBullet", parent=p_style, leftIndent=12, firstLineIndent=-8)

    story = [Paragraph(escape(title), title_style), Paragraph(f"生成日期：{datetime.date.today()} · 内部参考简报", meta_style)]
    for section_title, lines in parse_sections(report_text):
        if section_title:
            story.append(Paragraph(escape(normalize_line(section_title)), h_style))
        for block_type, content in detect_blocks(lines):
            if block_type == "table":
                header, rows = parse_markdown_table(content)
                if header:
                    data = [[Paragraph(escape(clean_cell_text(c)), p_style) for c in header]]
                    for row in rows[:12]:
                        row = row + [""] * (len(header) - len(row))
                        data.append([Paragraph(escape(clean_cell_text(c)), p_style) for c in row[:len(header)]])
                    table = Table(data, repeatRows=1)
                    table.setStyle(TableStyle([
                        ("BACKGROUND", (0,0), (-1,0), colors.HexColor("#eff4ff")),
                        ("GRID", (0,0), (-1,-1), 0.35, colors.HexColor("#dce6ff")),
                        ("VALIGN", (0,0), (-1,-1), "TOP"),
                        ("LEFTPADDING", (0,0), (-1,-1), 5),
                        ("RIGHTPADDING", (0,0), (-1,-1), 5),
                    ]))
                    story.append(table)
                    story.append(Spacer(1, 8))
            elif block_type in {"ul", "ol"}:
                for idx, item in enumerate(content, start=1):
                    prefix = "•" if block_type == "ul" else f"{idx}."
                    story.append(Paragraph(f"{prefix} {escape(clean_cell_text(item))}", bullet_style))
            elif block_type == "p":
                for para in content:
                    story.append(Paragraph(escape(clean_cell_text(para)), p_style))
    doc.build(story)
    return bio.getvalue()


def create_pptx_report(report_text: str, title: str) -> bytes:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    def add_title(slide, text):
        box = slide.shapes.add_textbox(Inches(0.65), Inches(0.38), Inches(12.0), Inches(0.55))
        p = box.text_frame.paragraphs[0]
        p.text = text
        p.font.name = "Microsoft YaHei"
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = RGBColor(35, 54, 111)
        return box

    def add_footer(slide):
        box = slide.shapes.add_textbox(Inches(0.65), Inches(7.03), Inches(12.0), Inches(0.22))
        p = box.text_frame.paragraphs[0]
        p.text = "Signal Studio · Evidence-based product brief"
        p.font.name = "Microsoft YaHei"
        p.font.size = Pt(8)
        p.font.color.rgb = RGBColor(122, 135, 155)
        p.alignment = PP_ALIGN.RIGHT

    def add_bullets(slide, items, x=0.85, y=1.25, w=11.7, h=5.6, font_size=16):
        box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = box.text_frame
        tf.clear()
        for idx, item in enumerate(items):
            p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
            p.text = item[:180]
            p.font.name = "Microsoft YaHei"
            p.font.size = Pt(font_size)
            p.font.color.rgb = RGBColor(34, 49, 79)
            p.level = 0
        return box

    # Cover
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = RGBColor(248, 251, 255)
    box = slide.shapes.add_textbox(Inches(0.8), Inches(1.7), Inches(11.8), Inches(1.6))
    p = box.text_frame.paragraphs[0]
    p.text = title
    p.font.name = "Microsoft YaHei"
    p.font.size = Pt(34)
    p.font.bold = True
    p.font.color.rgb = RGBColor(22, 32, 51)
    sub = slide.shapes.add_textbox(Inches(0.82), Inches(3.35), Inches(10.5), Inches(0.45))
    sp = sub.text_frame.paragraphs[0]
    sp.text = f"生成日期：{datetime.date.today()} · 内部汇报版"
    sp.font.name = "Microsoft YaHei"
    sp.font.size = Pt(13)
    sp.font.color.rgb = RGBColor(102, 116, 138)
    add_footer(slide)

    sections = parse_sections(report_text)
    summary = extract_summary_points(sections, max_items=5)
    if summary:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_title(slide, "老板先看")
        add_bullets(slide, [f"• {s}" for s in summary], font_size=17)
        add_footer(slide)

    # One slide per useful section, trimmed
    for section_title, lines in sections[1:8]:
        blocks = detect_blocks(lines)
        items = []
        for block_type, content in blocks:
            if block_type == "table":
                header, rows = parse_markdown_table(content)
                for row in rows[:5]:
                    items.append("｜".join(clean_cell_text(c) for c in row[:3]))
            elif block_type in {"ul", "ol"}:
                items.extend([clean_cell_text(i) for i in content[:6]])
            elif block_type == "p":
                items.extend([clean_cell_text(i) for i in content[:4]])
            if len(items) >= 7:
                break
        if not items:
            continue
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_title(slide, normalize_line(section_title)[:26])
        add_bullets(slide, [f"• {i}" for i in items[:7]], font_size=14)
        add_footer(slide)

    bio = BytesIO()
    prs.save(bio)
    return bio.getvalue()


def extract_boss_summary(report_text: str, title: str) -> str:
    sections = parse_sections(report_text)
    points = extract_summary_points(sections, max_items=5)
    if not points:
        # Fallback: first few non-table lines
        points = []
        for _, lines in sections[:2]:
            for block_type, content in detect_blocks(lines):
                if block_type in {"ul", "ol", "p"}:
                    points.extend(content)
                if len(points) >= 5:
                    break
            if len(points) >= 5:
                break
    body = "\n".join([f"{idx}. {clean_cell_text(p)}" for idx, p in enumerate(points[:5], start=1)])
    return f"【{title}】\n{body}"


def render_report_preview(report_text: str, title: str):
    st.markdown(build_report_body_html(report_text, title), unsafe_allow_html=True)

    st.markdown("#### 导出与复用")
    html_doc = build_download_html(report_text, title).encode("utf-8")
    base_name = sanitize_filename(title)
    export_cols = st.columns(5)
    with export_cols[0]:
        st.download_button("HTML / 可转 PDF", html_doc, file_name=base_name + ".html", mime="text/html", use_container_width=True)
    with export_cols[1]:
        st.download_button("Word", create_docx_report(report_text, title), file_name=base_name + ".docx", mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document", use_container_width=True)
    with export_cols[2]:
        st.download_button("PDF", create_pdf_report(report_text, title), file_name=base_name + ".pdf", mime="application/pdf", use_container_width=True)
    with export_cols[3]:
        st.download_button("PPT", create_pptx_report(report_text, title), file_name=base_name + ".pptx", mime="application/vnd.openxmlformats-officedocument.presentationml.presentation", use_container_width=True)
    with export_cols[4]:
        st.download_button("Markdown", report_text.encode("utf-8"), file_name=base_name + ".md", mime="text/markdown", use_container_width=True)

    with st.expander("老板版摘要 / 可复制", expanded=True):
        st.text_area("复制这段发微信/飞书/邮件", value=extract_boss_summary(report_text, title), height=160)

    with st.expander("查看原始 Markdown 报告", expanded=False):
        st.code(report_text, language="markdown")


# =========================
# Sidebar controls
# =========================
with st.sidebar:
    st.markdown("## Signal Studio")
    st.caption("PRODUCT INTELLIGENCE")

    api_key = st.text_input("API Key", type="password", placeholder="请输入你的 API Key")
    provider_name = st.selectbox("服务商预设", list(API_PRESETS.keys()), index=0)
    provider_info = API_PRESETS[provider_name]
    api_base = st.selectbox("接口地址", [provider_info["base"]], index=0)
    model_name = st.selectbox("模型代号", provider_info["models"], index=0)

    st.divider()
    mode = st.radio("任务模式", ["单品深度研判", "竞品对比攻防"])
    product_type = st.selectbox("产品类型", ["手机/消费电子", "汽车/智能座舱", "家电/IOT", "软件/App", "旅游/酒店产品", "其他"])
    focus = st.multiselect(
        "分析重点",
        ["用户痛点", "卖点感知", "价格价值感", "品牌认知", "系统/软件体验", "影像/性能", "外观设计", "渠道/服务", "传播话术"],
        default=["用户痛点", "卖点感知", "价格价值感", "传播话术"],
    )
    st.markdown("<div class='micro-tip'>配置栏默认收起。接口和模型均为预设选择；如需图片分析，请选择支持视觉输入的模型。</div>", unsafe_allow_html=True)


# =========================
# Main UI
# =========================
render_header()
render_search_dock()

if mode == "单品深度研判":
    with st.container(border=True):
        st.markdown("### 单品深度研判")
        st.caption("分析某一款产品的真实用户反馈、核心卖点感知、主要痛点与下一步传播建议。")

        product_name = st.text_input("产品名称", key="single_product_name", placeholder="例如：vivo X Fold5、OPPO Find N6、iPhone 17")
        if product_name.strip():
            default_query = compose_query(product_name, "全网口碑")
            with st.expander("快捷取证入口", expanded=False):
                render_link_buttons(build_search_links(default_query, ["UGC 社媒", "视频测评", "搜索引擎"])[:8], columns_per_row=4)
        else:
            st.info("先输入产品名称，或在上方搜索坞输入后点击“同步到单品分析”。")

        render_evidence_helper(product_name)

        user_input = st.text_area(
            "真实语料",
            height=300,
            placeholder="建议按平台分段粘贴，例如：\n【小红书】……\n【微博】……\n【B站】……\n【京东】……",
        )
        single_images = st.file_uploader(
            "图片证据（可选）",
            type=["png", "jpg", "jpeg", "webp"],
            accept_multiple_files=True,
            key="single_images",
            help="可上传评论截图、测评截图、产品图、海报、表格截图等。需选择支持视觉输入的模型。",
        )
        render_uploaded_images(single_images, "单品图片证据")

        if st.button("生成单品研判报告", type="primary", use_container_width=True):
            if is_ready_to_analyze_multimodal(api_key, [product_name], [(product_name or "单品", user_input, single_images)]):
                with st.spinner("正在基于真实语料和图片证据生成研判报告..."):
                    prompt = build_single_prompt(product_name.strip(), product_type, focus, user_input.strip() or "暂无文字语料，仅使用上传图片证据。")
                    report = analyze_with_llm(prompt, api_key, model_name, api_base, image_groups=[(product_name.strip(), single_images)])
                    title = f"{product_name.strip()}_单品研判简报"
                    st.markdown("### 报告预览")
                    render_report_preview(report, title)

else:
    with st.container(border=True):
        st.markdown("### 竞品对比攻防")
        st.caption("把本品与竞品的 UGC 反馈放在同一张策略桌上，判断防守点、进攻点与可用传播话术。")

        col1, col2 = st.columns(2, gap="large")
        with col1:
            main_product = st.text_input("本品名称", key="main_product", placeholder="例如：vivo X Fold5")
            if main_product.strip():
                with st.expander("本品快捷入口", expanded=False):
                    q = compose_query(main_product, "全网口碑")
                    render_link_buttons(build_search_links(q, ["UGC 社媒", "视频测评"])[:6], columns_per_row=3)
            main_input = st.text_area("本品真实语料", height=300, placeholder="粘贴本品评论/测评/用户反馈...")
            main_images = st.file_uploader(
                "本品图片证据（可选）",
                type=["png", "jpg", "jpeg", "webp"],
                accept_multiple_files=True,
                key="main_images",
            )
            render_uploaded_images(main_images, "本品图片证据")

        with col2:
            competitor_product = st.text_input("竞品名称", key="competitor_product", placeholder="例如：OPPO Find N6 / 华为 Mate X 系列")
            if competitor_product.strip():
                with st.expander("竞品快捷入口", expanded=False):
                    q = compose_query(competitor_product, "全网口碑")
                    render_link_buttons(build_search_links(q, ["UGC 社媒", "视频测评"])[:6], columns_per_row=3)
            competitor_input = st.text_area("竞品真实语料", height=300, placeholder="粘贴竞品评论/测评/用户反馈...")
            competitor_images = st.file_uploader(
                "竞品图片证据（可选）",
                type=["png", "jpg", "jpeg", "webp"],
                accept_multiple_files=True,
                key="competitor_images",
            )
            render_uploaded_images(competitor_images, "竞品图片证据")

        render_evidence_helper(main_product or competitor_product)

        if st.button("生成竞品攻防报告", type="primary", use_container_width=True):
            evidence_groups = [(main_product or "本品", main_input, main_images), (competitor_product or "竞品", competitor_input, competitor_images)]
            if is_ready_to_analyze_multimodal(api_key, [main_product, competitor_product], evidence_groups):
                with st.spinner("正在基于真实语料和图片证据构建竞品攻防报告..."):
                    prompt = build_compare_prompt(
                        main_product.strip(),
                        competitor_product.strip(),
                        product_type,
                        focus,
                        main_input.strip() or "暂无本品文字语料，仅使用本品上传图片证据。",
                        competitor_input.strip() or "暂无竞品文字语料，仅使用竞品上传图片证据。",
                    )
                    report = analyze_with_llm(
                        prompt,
                        api_key,
                        model_name,
                        api_base,
                        image_groups=[(f"本品：{main_product.strip()}", main_images), (f"竞品：{competitor_product.strip()}", competitor_images)],
                    )
                    title = f"{main_product.strip()}_vs_{competitor_product.strip()}_竞品攻防简报"
                    st.markdown("### 报告预览")
                    render_report_preview(report, title)
